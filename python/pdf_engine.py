"""PDF operations backend — shared by PyQt app and Electron CLI."""

import base64
import os
import shutil
import subprocess
import sys
import tempfile

try:
    import fitz
except ImportError as exc:
    raise ImportError("pymupdf is required: pip install pymupdf") from exc

try:
    from pypdf import PdfReader, PdfWriter
except ImportError as exc:
    raise ImportError("pypdf is required: pip install pypdf") from exc


def get_page_count(pdf_path):
    try:
        doc = fitz.open(pdf_path)
        n = doc.page_count
        doc.close()
        return n
    except Exception:
        return 0


def get_pdf_preview(pdf_path, dpi=96):
    """First-page thumbnail + page count for merge preview."""
    try:
        doc = fitz.open(pdf_path)
        pages = doc.page_count
        if pages == 0:
            doc.close()
            return False, 0, "", 0, 0
        page = doc[0]
        mat = fitz.Matrix(dpi / 72, dpi / 72)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        png_bytes = pix.tobytes("png")
        doc.close()
        return (
            True,
            pages,
            base64.b64encode(png_bytes).decode("ascii"),
            pix.width,
            pix.height,
        )
    except Exception as e:
        print(f"get_pdf_preview error: {e}")
        return False, 0, "", 0, 0


def render_page_png(
    pdf_path, page_index, dpi=150, zoom=1.0, highlights=None, active_hi=-1, rotation=0
):
    try:
        doc = fitz.open(pdf_path)
        page = doc[page_index]
        scale = (dpi / 72) * zoom

        if highlights:
            for i, r in enumerate(highlights):
                rect = fitz.Rect(r["x0"], r["y0"], r["x1"], r["y1"])
                is_active = i == active_hi
                shape = page.new_shape()
                shape.draw_rect(fitz.Rect(rect.x0, rect.y0 - 1, rect.x1, rect.y1 + 1))
                shape.finish(
                    fill=(1.0, 0.55, 0.0) if is_active else (1.0, 0.95, 0.0),
                    color=(0.8, 0.4, 0.0) if is_active else (0.85, 0.70, 0.0),
                    width=0.5,
                    fill_opacity=0.55 if is_active else 0.38,
                )
                shape.commit()

        mat = fitz.Matrix(scale, scale)
        if rotation:
            mat = mat.prerotate(int(rotation))
        pix = page.get_pixmap(matrix=mat, alpha=False)
        png_bytes = pix.tobytes("png")
        doc.close()
        return True, base64.b64encode(png_bytes).decode("ascii"), pix.width, pix.height
    except Exception as e:
        print(f"render_page_png error: {e}")
        return False, "", 0, 0


def reorder_pages(inp, out_path, order, cb=None):
    try:
        src = fitz.open(inp)
        out = fitz.open()
        total = len(order)
        for i, pg_i in enumerate(order):
            if cb:
                cb(i + 1, total, f"Arranging page {i + 1} of {total}…")
            out.insert_pdf(src, from_page=pg_i, to_page=pg_i)
        if cb:
            cb(total, total, "Saving…")
        out.save(out_path, deflate=True, garbage=4)
        src.close()
        out.close()
        return True
    except Exception as e:
        print(f"reorder_pages error: {e}")
        return False


def delete_pages(inp, out_path, pages_to_delete, cb=None):
    delete_set = set(pages_to_delete)
    try:
        src = fitz.open(inp)
        order = [i for i in range(src.page_count) if i not in delete_set]
        src.close()
        if not order:
            return False
        return reorder_pages(inp, out_path, order, cb)
    except Exception as e:
        print(f"delete_pages error: {e}")
        return False


def duplicate_pages(inp, out_path, page_indices, cb=None):
    try:
        src = fitz.open(inp)
        out = fitz.open()
        dup_set = set(page_indices)
        total = src.page_count + len(dup_set)
        step = [0]
        for i in range(src.page_count):
            if cb:
                step[0] += 1
                cb(step[0], total, f"Processing page {i + 1}…")
            out.insert_pdf(src, from_page=i, to_page=i)
            if i in dup_set:
                if cb:
                    step[0] += 1
                    cb(step[0], total, f"Duplicating page {i + 1}…")
                out.insert_pdf(src, from_page=i, to_page=i)
        if cb:
            cb(total, total, "Saving…")
        out.save(out_path, deflate=True, garbage=4)
        src.close()
        out.close()
        return True
    except Exception as e:
        print(f"duplicate_pages error: {e}")
        return False


def search_pdf(pdf_path, query):
    results = []
    if not query or not query.strip():
        return results
    try:
        doc = fitz.open(pdf_path)
        for i in range(doc.page_count):
            page = doc[i]
            for rect in page.search_for(query):
                results.append({
                    "page": i,
                    "x0": rect.x0,
                    "y0": rect.y0,
                    "x1": rect.x1,
                    "y1": rect.y1,
                })
        doc.close()
    except Exception as e:
        print(f"search_pdf error: {e}")
    return results


def extract_page_text(pdf_path, page_index):
    try:
        doc = fitz.open(pdf_path)
        txt = doc[page_index].get_text()
        doc.close()
        return txt or ""
    except Exception:
        return ""


def extract_all_text(pdf_path):
    try:
        doc = fitz.open(pdf_path)
        parts = [doc[i].get_text() or "" for i in range(doc.page_count)]
        doc.close()
        return "\n\n".join(parts)
    except Exception:
        return ""


def extract_pdf_for_ai(pdf_path):
    """Extract text + metadata for the AI assistant."""
    try:
        doc = fitz.open(pdf_path)
        parts = [doc[i].get_text() or "" for i in range(doc.page_count)]
        meta = doc.metadata or {}
        numpages = doc.page_count
        doc.close()
        return {
            "text": "\n\n".join(parts),
            "numpages": numpages,
            "info": {
                "Title": meta.get("title") or "",
                "Author": meta.get("author") or "",
                "Creator": meta.get("creator") or "",
                "Producer": meta.get("producer") or "",
                "CreationDate": meta.get("creationDate") or "",
                "ModDate": meta.get("modDate") or "",
            },
        }
    except Exception as e:
        print(f"extract_pdf_for_ai error: {e}")
        return {"text": "", "numpages": 0, "info": {}}


def _render_to_bytes(fitz_doc, q, dpi=150):
    out = fitz.open()
    scale = dpi / 72
    mat = fitz.Matrix(scale, scale)
    for page in fitz_doc:
        pix = page.get_pixmap(matrix=mat, alpha=False)
        jpeg_bytes = pix.tobytes(output="jpeg", jpg_quality=q)
        img_pdf_buf = fitz.open("jpeg", jpeg_bytes).convert_to_pdf()
        img_doc = fitz.open("pdf", img_pdf_buf)
        out.insert_pdf(img_doc)
        img_doc.close()
    buf = out.tobytes(deflate=True)
    out.close()
    return buf


def compress_to_target_kb(inp, out_path, target_kb, cb=None):
    tb = int(target_kb * 1024)
    if os.path.getsize(inp) <= tb:
        shutil.copy2(inp, out_path)
        if cb:
            cb(100, 100, "Already within target.")
        return True, os.path.getsize(inp) / 1024
    src = fitz.open(inp)
    best = None
    dpis = [150, 120, 96, 72, 48, 30]
    total = len(dpis) * 8
    step = [0]
    for dpi in dpis:
        lo, hi, found = 10, 85, False
        for _ in range(8):
            mq = (lo + hi) // 2
            step[0] += 1
            if cb:
                cb(step[0], total, f"Trying {dpi} DPI  q{mq}…")
            buf = _render_to_bytes(src, mq, dpi)
            if len(buf) <= tb:
                best = buf
                lo = mq + 1
                found = True
            else:
                hi = mq - 1
        if found:
            break
    src.close()
    if best is None:
        s2 = fitz.open(inp)
        best = _render_to_bytes(s2, 10, 30)
        s2.close()
    if cb:
        cb(total, total, "Saving…")
    with open(out_path, "wb") as f:
        f.write(best)
    return True, len(best) / 1024


def resize_pdf_pages(inp, out_path, w_mm, h_mm, cb=None):
    try:
        pts_per_mm = 72.0 / 25.4
        nw = w_mm * pts_per_mm
        nh = h_mm * pts_per_mm
        src = fitz.open(inp)
        out = fitz.open()
        tot = src.page_count
        for i in range(tot):
            if cb:
                cb(i + 1, tot, f"Resizing page {i + 1} of {tot}…")
            pg = src[i]
            src_w = pg.rect.width
            src_h = pg.rect.height
            mat = fitz.Matrix(150 / 72, 150 / 72)
            pix = pg.get_pixmap(matrix=mat, alpha=False)
            img_bytes = pix.tobytes("jpeg", jpg_quality=92)
            img_pdf_buf = fitz.open("jpeg", img_bytes).convert_to_pdf()
            img_doc = fitz.open("pdf", img_pdf_buf)
            scale = min(nw / src_w, nh / src_h)
            dst_w = src_w * scale
            dst_h = src_h * scale
            x0 = (nw - dst_w) / 2.0
            y0 = (nh - dst_h) / 2.0
            new_pg = out.new_page(width=nw, height=nh)
            new_pg.show_pdf_page(
                fitz.Rect(x0, y0, x0 + dst_w, y0 + dst_h), img_doc, 0
            )
            img_doc.close()
        if cb:
            cb(tot, tot, "Saving…")
        out.save(out_path, deflate=True, garbage=4, clean=True)
        src.close()
        out.close()
        return True
    except Exception as e:
        print(f"resize_pdf_pages error: {e}")
        return False


def compress_standard(inp, out_path, level="balanced", cb=None):
    try:
        src = fitz.open(inp)
        if cb:
            cb(10, 100, "Compressing…")
        if level == "light":
            src.save(out_path, deflate=True, garbage=4, clean=True)
        else:
            dpi = 120 if level == "balanced" else 72
            q = 85 if level == "balanced" else 60
            buf = _render_to_bytes(src, q, dpi)
            if cb:
                cb(90, 100, "Saving…")
            with open(out_path, "wb") as f:
                f.write(buf)
        src.close()
        if cb:
            cb(100, 100, "Done.")
        return True
    except Exception as e:
        print(e)
        return False


def split_pdf_range(inp, out_path, start_page, end_page, cb=None):
    try:
        src = fitz.open(inp)
        new_pdf = fitz.open()
        total = end_page - start_page + 1
        for i, pg_i in enumerate(range(start_page, end_page + 1)):
            if cb:
                cb(i + 1, total, f"Extracting page {pg_i + 1}...")
            new_pdf.insert_pdf(src, from_page=pg_i, to_page=pg_i)
        if cb:
            cb(total, total, "Saving...")
        new_pdf.save(out_path, deflate=True, garbage=4)
        new_pdf.close()
        src.close()
        return True
    except Exception as e:
        print(f"split_pdf_range error: {e}")
        return False


def split_pdf_individual(inp, out_dir, base_name, cb=None):
    try:
        src = fitz.open(inp)
        total = src.page_count
        paths = []
        for i in range(total):
            if cb:
                cb(i + 1, total, f"Saving page {i + 1} of {total}...")
            one = fitz.open()
            one.insert_pdf(src, from_page=i, to_page=i)
            p = os.path.join(out_dir, f"{base_name}_p{i+1:03d}.pdf")
            one.save(p, deflate=True, garbage=4)
            one.close()
            paths.append(p)
        src.close()
        if cb:
            cb(total, total, "Done.")
        return True, paths
    except Exception as e:
        print(f"split_pdf_individual error: {e}")
        return False, []


def rotate_pdf_pages_map(inp, out_path, page_degrees, cb=None):
    try:
        src = fitz.open(inp)
        total = len(page_degrees) or 1
        for n, (pg_i, deg) in enumerate(page_degrees.items()):
            idx = int(pg_i)
            rot = int(deg) % 360
            if cb:
                cb(n + 1, total, f"Rotating page {idx + 1}…")
            if 0 <= idx < src.page_count and rot:
                pg = src[idx]
                pg.set_rotation((pg.rotation + rot) % 360)
        if cb:
            cb(total, total, "Saving…")
        fd, tmp_path = tempfile.mkstemp(suffix=".pdf")
        os.close(fd)
        src.save(tmp_path, deflate=True, garbage=4, clean=True, incremental=False)
        src.close()
        shutil.move(tmp_path, out_path)
        return True
    except Exception as e:
        print(f"rotate_pdf_pages_map error: {e}")
        return False


def get_text_spans(pdf_path, page_index):
    spans = []
    try:
        doc = fitz.open(pdf_path)
        if page_index < 0 or page_index >= doc.page_count:
            doc.close()
            return spans
        page = doc[page_index]
        for block in page.get_text("dict").get("blocks", []):
            if block.get("type") != 0:
                continue
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    txt = span.get("text", "")
                    if not txt.strip():
                        continue
                    bbox = span["bbox"]
                    spans.append({
                        "text": txt,
                        "x0": bbox[0],
                        "y0": bbox[1],
                        "x1": bbox[2],
                        "y1": bbox[3],
                        "size": span.get("size", 11),
                    })
        doc.close()
    except Exception as e:
        print(f"get_text_spans error: {e}")
    return spans


def tesseract_available():
    return bool(shutil.which("tesseract"))


def rotate_pdf_pages(inp, out_path, degrees, page_indices=None, cb=None):
    tmp_path = None
    try:
        src = fitz.open(inp)
        total = src.page_count
        targets = set(page_indices) if page_indices is not None else set(range(total))
        for i in range(total):
            if cb:
                cb(i + 1, total, f"Processing page {i + 1}...")
            if i in targets:
                pg = src[i]
                pg.set_rotation((pg.rotation + degrees) % 360)
        if cb:
            cb(total, total, "Saving...")
        fd, tmp_path = tempfile.mkstemp(suffix=".pdf")
        os.close(fd)
        src.save(tmp_path, deflate=True, garbage=4, clean=True, incremental=False)
        src.close()
        shutil.move(tmp_path, out_path)
        tmp_path = None
        return True
    except Exception as e:
        print(f"rotate_pdf_pages error: {e}")
        return False
    finally:
        if tmp_path and os.path.exists(tmp_path):
            try:
                os.remove(tmp_path)
            except OSError:
                pass


def add_watermark(inp, out_path, text, opacity=0.25, cb=None):
    try:
        src = fitz.open(inp)
        total = src.page_count
        gray = max(0.05, min(0.80, 1.0 - opacity))
        color = (gray, gray, gray)
        for i in range(total):
            if cb:
                cb(i + 1, total, f"Watermarking page {i + 1}...")
            pg = src[i]
            r = pg.rect
            fontsize = max(18, min(80, int(r.width / max(len(text), 1) * 1.8)))
            positions = [
                fitz.Point(r.width * 0.15, r.height * 0.45),
                fitz.Point(r.width * 0.30, r.height * 0.70),
            ]
            mat45 = fitz.Matrix(45)
            for pt in positions:
                pg.insert_text(
                    pt, text, fontsize=fontsize, color=color,
                    rotate=0, morph=(pt, mat45), overlay=True,
                )
        if cb:
            cb(total, total, "Saving...")
        pdf_bytes = src.tobytes(deflate=True, garbage=4)
        src.close()
        with open(out_path, "wb") as f:
            f.write(pdf_bytes)
        return True
    except Exception as e:
        print(f"add_watermark error: {e}")
        return False


def password_protect_pdf(inp, out_path, user_pw, owner_pw=None, cb=None):
    try:
        if cb:
            cb(10, 100, "Opening PDF...")
        src = fitz.open(inp)
        if not src.is_pdf:
            src.close()
            return False

        def _const(name, fallback):
            return getattr(fitz, name, fallback)

        enc = _const("PDF_ENCRYPT_AES_256", 6)
        perm = (
            _const("PDF_PERM_PRINT", 4)
            | _const("PDF_PERM_COPY", 16)
            | _const("PDF_PERM_PRINT_HQ", 2048)
        )
        if cb:
            cb(50, 100, "Encrypting...")
        src.save(
            out_path,
            encryption=enc,
            user_pw=user_pw,
            owner_pw=owner_pw or user_pw,
            permissions=perm,
            garbage=4,
            deflate=True,
        )
        src.close()
        if cb:
            cb(100, 100, "Done.")
        return True
    except Exception as e:
        print(f"password_protect_pdf error: {e}")
        return False


def merge_pdfs(paths, out_path, cb=None):
    try:
        w = PdfWriter()
        total = len(paths)
        for i, p in enumerate(paths):
            if cb:
                cb(i + 1, total, f"Merging file {i + 1} of {total}...")
            for pg in PdfReader(p).pages:
                w.add_page(pg)
        if cb:
            cb(total, total, "Saving...")
        with open(out_path, "wb") as f:
            w.write(f)
        return True
    except Exception as e:
        print(f"merge_pdfs error: {e}")
        return False


def get_tts_voices():
    try:
        import pyttsx3
        e = pyttsx3.init()
        voices = e.getProperty("voices") or []
        result = [{"name": v.name, "id": v.id} for v in voices]
        e.stop()
        return result
    except Exception:
        return []


def pdf_needs_ocr(pdf_path, sample_pages=5):
    try:
        doc = fitz.open(pdf_path)
        chars = 0
        for i in range(min(sample_pages, doc.page_count)):
            chars += len(doc[i].get_text().strip())
        doc.close()
        return chars < 40
    except Exception:
        return False


def ocr_pdf(inp, out_path, language="eng", cb=None):
    """Add searchable text layer to scanned pages using Tesseract OCR."""
    if not tesseract_available():
        print("OCR error: Tesseract not found on PATH")
        return False
    try:
        src = fitz.open(inp)
        out = fitz.open()
        total = src.page_count
        ocr_count = 0

        for i in range(total):
            if cb:
                cb(i + 1, total, f"OCR page {i + 1} of {total}…")
            pg = src[i]
            new_pg = out.new_page(width=pg.rect.width, height=pg.rect.height)
            new_pg.insert_image(pg.rect, pixmap=pg.get_pixmap(dpi=200))

            try:
                tp = pg.get_textpage_ocr(language=language, dpi=300, full=True)
                blocks = tp.extractDICT().get("blocks", [])
                page_chars = 0
                for block in blocks:
                    if block.get("type") != 0:
                        continue
                    for line in block.get("lines", []):
                        for span in line.get("spans", []):
                            txt = span.get("text", "").strip()
                            if not txt:
                                continue
                            bbox = span["bbox"]
                            fontsize = max(6, span.get("size", 10))
                            new_pg.insert_text(
                                fitz.Point(bbox[0], bbox[3]),
                                txt,
                                fontsize=fontsize,
                                fontname="helv",
                                color=(0, 0, 0),
                                render_mode=3,
                            )
                            page_chars += len(txt)
                if page_chars:
                    ocr_count += 1
            except Exception as ex:
                print(f"OCR page {i + 1} error: {ex}")

        if cb:
            cb(total, total, "Saving searchable PDF…")
        out.save(out_path, deflate=True, garbage=4)
        src.close()
        out.close()
        return ocr_count > 0 or total > 0
    except Exception as e:
        print(f"ocr_pdf error: {e}")
        import traceback
        traceback.print_exc()
        return False


def apply_edits(inp, out_path, edits, cb=None):
    """Apply text, highlight, and rectangle edits to a PDF."""
    try:
        src = fitz.open(inp)
        total = len(edits) or 1
        for idx, ed in enumerate(edits):
            if cb:
                cb(idx + 1, total, f"Applying edit {idx + 1}…")
            page_i = ed.get("page", 0)
            if page_i < 0 or page_i >= src.page_count:
                continue
            pg = src[page_i]
            etype = ed.get("type")
            if etype == "text":
                pg.insert_text(
                    fitz.Point(ed["x"], ed["y"]),
                    ed["text"],
                    fontsize=ed.get("size", 12),
                    color=tuple(ed.get("color", [0, 0, 0])),
                )
            elif etype == "highlight":
                rect = fitz.Rect(ed["x0"], ed["y0"], ed["x1"], ed["y1"])
                annot = pg.add_highlight_annot(rect)
                annot.set_colors(stroke=None)
                annot.update()
            elif etype == "rect":
                rect = fitz.Rect(ed["x0"], ed["y0"], ed["x1"], ed["y1"])
                shape = pg.new_shape()
                shape.draw_rect(rect)
                shape.finish(
                    width=ed.get("width", 1.5),
                    color=tuple(ed.get("color", [0.23, 0.39, 0.9])),
                )
                shape.commit()
            elif etype == "replace_text":
                rect = fitz.Rect(ed["x0"], ed["y0"], ed["x1"], ed["y1"])
                pg.add_redact_annot(rect, fill=(1, 1, 1))
                pg.apply_redactions()
                pg.insert_text(
                    fitz.Point(ed["x0"], ed["y1"] - 1),
                    ed["text"],
                    fontsize=ed.get("size", 11),
                    fontname="helv",
                    color=tuple(ed.get("color", [0, 0, 0])),
                )
        if cb:
            cb(total, total, "Saving…")
        src.save(out_path, deflate=True, garbage=4)
        src.close()
        return True
    except Exception as e:
        print(f"apply_edits error: {e}")
        import traceback
        traceback.print_exc()
        return False


def save_tts_audio(text, out_path, voice_id=None, rate=175):
    try:
        import pyttsx3
        e = pyttsx3.init()
        e.setProperty("rate", rate)
        if voice_id:
            e.setProperty("voice", voice_id)
        e.save_to_file(text, out_path)
        e.runAndWait()
        return True
    except Exception as e:
        print(f"save_tts_audio error: {e}")
        return False


# ── File conversion ──

def _find_soffice():
    candidates = [
        "soffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    for c in candidates:
        if shutil.which(c):
            return c
        if os.path.isfile(c):
            return c
    return None


def _libreoffice_convert(inp, out_path, target_fmt, cb=None):
    soffice = _find_soffice()
    if not soffice:
        return False
    inp = os.path.abspath(inp)
    out_path = os.path.abspath(out_path)
    out_dir = os.path.dirname(out_path) or os.getcwd()
    os.makedirs(out_dir, exist_ok=True)
    if cb:
        cb(1, 2, f"Converting with LibreOffice…")
    result = subprocess.run(
        [soffice, "--headless", "--convert-to", target_fmt, "--outdir", out_dir, inp],
        capture_output=True,
        text=True,
        timeout=180,
        cwd=out_dir,
    )
    base = os.path.splitext(os.path.basename(inp))[0]
    generated = os.path.join(out_dir, f"{base}.{target_fmt}")
    if os.path.isfile(generated):
        if os.path.abspath(generated) != os.path.abspath(out_path):
            if os.path.isfile(out_path):
                os.remove(out_path)
            os.rename(generated, out_path)
        return os.path.isfile(out_path)
    if result.returncode == 0 and os.path.isfile(out_path):
        return True
    print(f"LibreOffice convert failed: {(result.stderr or result.stdout or '')[:400]}")
    return False


def _docx_extract_lines(docx):
    from docx import Document

    doc = Document(docx)
    lines = []
    for para in doc.paragraphs:
        if para.text.strip():
            lines.append(para.text.strip())
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                lines.append(" | ".join(cells))
    return lines


def _docx_to_pdf_fallback(inp, out_path, cb=None):
    if cb:
        cb(1, 3, "Converting Word document…")
    try:
        import mammoth

        with open(inp, "rb") as f:
            result = mammoth.convert_to_html(f)
        html = (
            '<div style="font-family: Arial, sans-serif; font-size: 11pt; line-height: 1.5;">'
            f"{result.value}</div>"
        )
        if cb:
            cb(2, 3, "Building PDF pages…")
        pdf = fitz.open()
        page = pdf.new_page(width=595, height=842)
        rect = fitz.Rect(40, 40, 555, 802)
        remainder = page.insert_htmlbox(rect, html)
        while remainder:
            page = pdf.new_page(width=595, height=842)
            rect = fitz.Rect(40, 40, 555, 802)
            remainder = page.insert_htmlbox(rect, remainder)
        pdf.save(out_path)
        pdf.close()
        if os.path.isfile(out_path) and os.path.getsize(out_path) > 500:
            return True
    except Exception as e:
        print(f"mammoth docx error: {e}")

    try:
        if cb:
            cb(2, 3, "Using text extraction fallback…")
        lines = _docx_extract_lines(inp)
        if not lines:
            return False
        pdf = fitz.open()
        page = pdf.new_page(width=595, height=842)
        y = 56
        for text in lines:
            if y > 780:
                page = pdf.new_page(width=595, height=842)
                y = 56
            page.insert_text((56, y), text, fontsize=11, fontname="helv")
            y += 16
        pdf.save(out_path)
        pdf.close()
        return os.path.isfile(out_path)
    except Exception as e:
        print(f"docx fallback error: {e}")
        return False


def _xlsx_to_pdf_fallback(inp, out_path, cb=None):
    try:
        from openpyxl import load_workbook

        if cb:
            cb(1, 2, "Converting Excel spreadsheet…")
        wb = load_workbook(inp, data_only=True)
        pdf = fitz.open()
        for sheet in wb.worksheets:
            page = pdf.new_page(width=842, height=595)
            y = 40
            page.insert_text((40, y), sheet.title, fontsize=14)
            y += 24
            for row in sheet.iter_rows(max_row=min(sheet.max_row or 0, 80), values_only=True):
                if y > 560:
                    page = pdf.new_page(width=842, height=595)
                    y = 40
                cells = ["" if v is None else str(v) for v in row]
                if not any(cells):
                    continue
                line = " | ".join(cells)[:220]
                page.insert_text((40, y), line, fontsize=9)
                y += 12
        pdf.save(out_path)
        pdf.close()
        return os.path.isfile(out_path)
    except Exception as e:
        print(f"xlsx fallback error: {e}")
        return False


def _pptx_to_pdf_fallback(inp, out_path, cb=None):
    try:
        from pptx import Presentation

        if cb:
            cb(1, 2, "Converting PowerPoint…")
        prs = Presentation(inp)
        pdf = fitz.open()
        for slide in prs.slides:
            page = pdf.new_page(width=960, height=540)
            y = 48
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    for line in shape.text.splitlines():
                        if y > 500:
                            break
                        page.insert_text((48, y), line[:180], fontsize=12)
                        y += 18
        pdf.save(out_path)
        pdf.close()
        return os.path.isfile(out_path)
    except Exception as e:
        print(f"pptx fallback error: {e}")
        return False


def office_to_pdf(inp, out_path, cb=None):
    ext = os.path.splitext(inp)[1].lower()
    office_exts = (".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx", ".odt", ".ods", ".odp", ".csv")
    if ext not in office_exts:
        return False, "Unsupported office format."

    if _libreoffice_convert(inp, out_path, "pdf", cb):
        return True, "Converted to PDF."

    if sys.platform == "win32" and ext in (".docx", ".doc"):
        try:
            from docx2pdf import convert

            if cb:
                cb(1, 2, "Converting with Microsoft Word…")
            convert(os.path.abspath(inp), os.path.abspath(out_path))
            if os.path.isfile(out_path):
                return True, "Converted to PDF."
        except Exception as e:
            print(f"docx2pdf error: {e}")

    if ext in (".docx", ".doc", ".odt") and _docx_to_pdf_fallback(inp, out_path, cb):
        return True, "Converted to PDF (basic layout)."
    if ext in (".xlsx", ".xls", ".ods", ".csv") and _xlsx_to_pdf_fallback(inp, out_path, cb):
        return True, "Converted to PDF (basic layout)."
    if ext in (".pptx", ".ppt", ".odp") and _pptx_to_pdf_fallback(inp, out_path, cb):
        return True, "Converted to PDF (basic layout)."

    return False, "Could not convert file. Install LibreOffice for best results."


def _pdf_to_xlsx_fallback(inp, out_path, cb=None):
    try:
        from openpyxl import Workbook

        if cb:
            cb(1, 2, "Exporting to Excel…")
        doc = fitz.open(inp)
        wb = Workbook()
        wb.remove(wb.active)
        for i in range(doc.page_count):
            ws = wb.create_sheet(title=f"Page {i + 1}"[:31])
            text = doc[i].get_text()
            for row_idx, line in enumerate(text.splitlines(), start=1):
                if line.strip():
                    ws.cell(row=row_idx, column=1, value=line.strip())
        doc.close()
        wb.save(out_path)
        return os.path.isfile(out_path)
    except Exception as e:
        print(f"pdf to xlsx fallback error: {e}")
        return False


def _pdf_to_pptx_fallback(inp, out_path, cb=None):
    try:
        from pptx import Presentation
        from pptx.util import Inches

        if cb:
            cb(1, 2, "Exporting to PowerPoint…")
        doc = fitz.open(inp)
        prs = Presentation()
        blank = prs.slide_layouts[6]
        total = doc.page_count
        for i in range(total):
            if cb:
                cb(i + 1, total, f"Slide {i + 1}…")
            pix = doc[i].get_pixmap(dpi=150)
            img_path = os.path.join(tempfile.gettempdir(), f"pdf_slide_{i}.png")
            pix.save(img_path)
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            try:
                os.remove(img_path)
            except OSError:
                pass
        doc.close()
        prs.save(out_path)
        return os.path.isfile(out_path)
    except Exception as e:
        print(f"pdf to pptx fallback error: {e}")
        return False


def pdf_to_office(inp, out_path, target_fmt, cb=None):
    target_fmt = target_fmt.lower().lstrip(".")

    if _libreoffice_convert(inp, out_path, target_fmt, cb):
        return True, f"Exported to {target_fmt.upper()}."

    if target_fmt == "docx":
        try:
            from pdf2docx import Converter

            if cb:
                cb(1, 2, "Converting to Word…")
            cv = Converter(inp)
            cv.convert(out_path)
            cv.close()
            if os.path.isfile(out_path):
                return True, "Exported to Word."
        except Exception as e:
            print(f"pdf2docx error: {e}")

    if target_fmt == "xlsx" and _pdf_to_xlsx_fallback(inp, out_path, cb):
        return True, "Exported to Excel (text per page)."

    if target_fmt == "pptx" and _pdf_to_pptx_fallback(inp, out_path, cb):
        return True, "Exported to PowerPoint (pages as slides)."

    return False, f"Could not export to {target_fmt.upper()}."


def _image_file_to_pdf_page(img_path):
    try:
        img = fitz.open(img_path)
        pdfbytes = img.convert_to_pdf()
        img.close()
        return fitz.open("pdf", pdfbytes)
    except Exception:
        from PIL import Image
        import io

        pil = Image.open(img_path).convert("RGB")
        buf = io.BytesIO()
        pil.save(buf, format="PNG")
        buf.seek(0)
        img = fitz.open(stream=buf.read(), filetype="png")
        pdfbytes = img.convert_to_pdf()
        img.close()
        return fitz.open("pdf", pdfbytes)


def images_to_pdf(image_paths, out_path, cb=None):
    try:
        if not image_paths:
            return False, "No images selected."
        doc = fitz.open()
        total = len(image_paths)
        for i, img_path in enumerate(image_paths):
            if not os.path.isfile(img_path):
                continue
            if cb:
                cb(i + 1, total, f"Adding image {i + 1}…")
            img_pdf = _image_file_to_pdf_page(img_path)
            doc.insert_pdf(img_pdf)
            img_pdf.close()
        if doc.page_count == 0:
            doc.close()
            return False, "No valid images found."
        if cb:
            cb(total, total, "Saving PDF…")
        os.makedirs(os.path.dirname(os.path.abspath(out_path)) or ".", exist_ok=True)
        doc.save(out_path, deflate=True)
        doc.close()
        return True, "Images combined into PDF."
    except Exception as e:
        print(f"images_to_pdf error: {e}")
        return False, str(e)


def pdf_to_images(inp, out_dir, image_fmt="png", dpi=150, cb=None):
    try:
        doc = fitz.open(inp)
        total = doc.page_count
        paths = []
        for i in range(total):
            if cb:
                cb(i + 1, total, f"Exporting page {i + 1}…")
            pix = doc[i].get_pixmap(dpi=dpi)
            out = os.path.join(out_dir, f"page_{i + 1:03d}.{image_fmt}")
            pix.save(out)
            paths.append(out)
        doc.close()
        return True, paths
    except Exception as e:
        print(f"pdf_to_images error: {e}")
        return False, []
